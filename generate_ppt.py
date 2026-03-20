from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title, content, notes=None, is_title_slide=False):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0] # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = content
    else:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        
        if isinstance(content, list):
            for i, item in enumerate(content):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = str(item)
                p.level = 0
        else:
            tf.text = str(content)
            
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

def create_presentation():
    prs = Presentation()
    
    # Slide 1
    add_slide(prs, 
        "BreathySync", 
        "AI-Powered Respiratory Health Monitoring Platform\n\n\"Breathe Smarter. Live Better.\"", 
        notes="Good [morning/afternoon]. I'm going to present BreathySync — an end-to-end respiratory health platform that uses your smartphone's microphone, real-time air quality data, and AI prediction to help 30 million asthma patients in India manage their health proactively. Let's dive in.",
        is_title_slide=True)
        
    # Slide 2
    add_slide(prs,
        "The Crisis Hidden in Every Breath",
        [
            "545 million people globally suffer from respiratory diseases (WHO).",
            "30 million asthma patients in India — 3rd leading cause of disability.",
            "40% of asthma hospitalizations are preventable with early intervention.",
            "Existing tools are fragmented: AQI apps, manual diaries, clinical spirometers — none connected.",
            "Nocturnal cough 24–48 hrs before attack = strongest early predictor — almost never monitored."
        ],
        notes="The core problem is that respiratory health management is reactive and fragmented. Patients only seek help after an attack. All the data that could predict that attack — voice changes, air quality spikes, sleep patterns — exists in silos. BreathySync connects these dots for the first time.")
        
    # Slide 3
    add_slide(prs,
        "BreathySync — 5 Pillars of Lung Health",
        [
            "1. Voice Biomarker Analysis: 6-sec voice → Lung Score (0–100)",
            "2. Environmental Risk Intelligence: Live AQI + Weather → Trigger alerts",
            "3. Digital Twin AI Prediction: Predicts exacerbation 24–48 hrs early",
            "4. Acoustic Sleep Radar: Passive nocturnal cough monitoring",
            "5. ABHA Integration: Government health record interoperability"
        ],
        notes="BreathySync is built on five pillars. Together they form a continuous monitoring loop — from sleep through the day and back — that no other platform offers. And critically, everything runs on a standard smartphone — no special hardware needed.")
        
    # Slide 4
    add_slide(prs,
        "System Architecture",
        [
            "Frontend: React 18 + TypeScript + Vite + Tailwind CSS",
            "Backend: FastAPI (Python) running audio DSP",
            "Database: Supabase (PostgreSQL) + Auth",
            "Data flow: Browser (MediaRecorder) -> API -> librosa analysis -> Supabase",
            "External Integrations: AQICN (Air), OpenWeatherMap, Gemini AI (Travel Chat)"
        ],
        notes="The architecture is cleanly separated into three layers. The React SPA handles all UI and user interactions. FastAPI handles all compute-heavy tasks like audio analysis. Supabase provides authentication, database, and real-time capabilities. All layers communicate via REST APIs secured with JWT tokens.")

    # Slide 5
    add_slide(prs,
        "Feature — Voice Check",
        [
            "Record 6 seconds of sustained 'Ahhh'.",
            "Browser converts WebM → WAV (client-side).",
            "Audio is uploaded to FastAPI for acoustic analysis.",
            "Extracts: Jitter, Shimmer, HNR, 39 MFCCs, Formants.",
            "Clinical validity: Maps DSP features to a 0–100 Lung Score.",
            "Replaces ₹50,000+ spirometer with a smartphone mic."
        ],
        notes="This is the crown jewel of BreathySync. Using your phone's microphone, we extract the same acoustic biomarkers that clinical voice labs measure. The YIN pitch tracking algorithm runs on our server in under 3 seconds. The result is a clinically-grounded score — not a random number.")

    # Slide 6
    add_slide(prs,
        "Feature — Digital Twin Predictor",
        [
            "Real-time exacerbation probability: 0–100% risk score.",
            "Integrates: Live AQI (30%), Voice Trend (50%), Trigger Reports (20%).",
            "Maps to GINA 2023 Action Plan tiers (Step 1–4).",
            "Example: AQI 180 + voice score declining = 82% risk (Escalate to Step 3).",
            "Proactive alert system instead of reactive treatment."
        ],
        notes="Think of the Digital Twin as an AI model of your lungs. It continuously evaluates your environmental conditions and recent health trajectory to estimate how close you are to an attack — and what you should do about it. This is the kind of clinical intelligence that used to require a pulmonologist visit.")

    # Slide 7
    add_slide(prs,
        "Feature — Sleep Mode & Trigger Map",
        [
            "Sleep Mode: Passive mic monitoring during sleep.",
            "Detects cough/wheeze frequency without storing audio (Privacy-by-design).",
            "Trigger Map: Interactive Leaflet map with live AQI circles.",
            "AI Travel Chat: 'Is Pune safe for my asthma this weekend?'",
            "Gemini AI generates risk assessment based on user profile."
        ],
        notes="Two features that work together. Sleep mode gives us the nocturnal data that predicts attacks 24 to 48 hours ahead. The Trigger Map contextualizes that data with where the patient is and where they're planning to go. The AI chat makes this accessible to non-technical users.")

    # Slide 8
    add_slide(prs,
        "Feature — Lung Gym & Gut Health",
        [
            "Dragon Breather: Microphone-controlled game (blow to fly).",
            "Converts boring breathing exercises into gamified engagement.",
            "Real-time Web Audio API volume detection.",
            "Gut-Lung Health: Tracks daily food log compliance.",
            "Recharts visualization correlates gut health with lung score."
        ],
        notes="Studies show that the biggest problem in asthma management isn't knowledge — it's adherence. Patients don't do their breathing exercises because they're boring. Dragon Breather makes clinical breathing exercises into a game. Similarly, the gut-lung tracker is based on published research linking microbiome health to asthma severity.")

    # Slide 9
    add_slide(prs,
        "Tech Stack — Frontend",
        [
            "React 18 & TypeScript: Type-safe, component-based UI.",
            "Vite: Lightning-fast build tool and development server.",
            "Framer Motion: Physics-based, GPU-accelerated micro-animations.",
            "Tailwind CSS: Zero-runtime utility classes for rapid styling.",
            "Leaflet.js: Open-source mapping (no API billing).",
            "Recharts: React-native data visualization.",
            "Supabase JS: Unified SDK for Auth and Database."
        ],
        notes="Every technology choice was deliberate. We chose Vite over Create React App for a dramatically faster development experience. Framer Motion over CSS animations for physics-based micro-interactions. Leaflet over Google Maps to avoid per-request billing. These aren't defaults — they're engineered decisions.")

    # Slide 10
    add_slide(prs,
        "Tech Stack — Backend & APIs",
        [
            "FastAPI (Python): Native async, auto Swagger docs.",
            "librosa (0.11): Industry-standard DSP for audio analysis.",
            "YIN Algorithm: Pitch tracking 10x faster than PYIN.",
            "NumPy & SciPy: Scientific and numerical computation.",
            "Supabase / PostgreSQL: RLS, Auth, REST APIs.",
            "External APIs: Google Gemini, AQICN, OpenWeatherMap."
        ],
        notes="The backend is pure Python, chosen specifically because librosa — the gold standard for audio feature extraction in machine learning research — is a Python library. FastAPI gives us automatic API documentation, strict Pydantic validation, and async I/O. The Gemini integration gives our chatbot genuine medical reasoning.")

    # Slide 11
    add_slide(prs,
        "Algorithms — Voice Analysis",
        [
            "YIN Pitch Tracking: O(N log N) frequency estimation.",
            "Jitter (Freq Perturbation): Healthy < 1.0%.",
            "Shimmer (Amp Perturbation): Healthy < 3.0%.",
            "HNR (Harmonic-to-Noise Ratio): Healthy > 20 dB.",
            "Deterministic lung score subtracts penalties for deviances."
        ],
        notes="These aren't arbitrary scores. Each metric is based on published clinical research. Jitter and Shimmer are standard markers in voice pathology labs. HNR measures voice quality degradation. Our rule-based model maps each deviation to a penalty, giving a transparent, explainable score — not a black-box output.")

    # Slide 12
    add_slide(prs,
        "Algorithms — Digital Twin",
        [
            "Weighted Heuristic Risk = (AQI×0.3) + (Voice Decline×0.5) + (Triggers×0.2).",
            "Maps to 4 clinical tiers: Low, Moderate, High, Critical.",
            "Transparent and explainable to clinicians.",
            "Future Roadmap: Transition to Random Forest trained on mPower dataset."
        ],
        notes="For the hackathon MVP, we use a well-designed heuristic model where all weights are grounded in clinical evidence — 50% weight on voice score because it's our hardest-to-fake signal. The roadmap calls for replacing this with a trained random forest model using public health datasets like the Parkinson's mPower study.")

    # Slide 13
    add_slide(prs,
        "Security Architecture",
        [
            "Authentication: Supabase Auth (OAuth 2.0).",
            "Access Control: PostgREST with Row-Level Security (RLS).",
            "Audio Privacy: Temporarily stored in memory, never saved to disk.",
            "Sleep Mode Privacy: Transmits only integer event counts, not raw audio.",
            "ABHA Privacy: Simulation only, no PII stored.",
            "Protection: React auto-escaping (XSS) and JWT headers (CSRF)."
        ],
        notes="Security and privacy weren't afterthoughts — they're baked in. Row-Level Security at the database means that even a server compromise can't expose other users' data. The privacy-by-design principle in Sleep Mode — transmitting only counts, never audio — is how we'd pass a real HIPAA audit.")

    # Slide 14
    add_slide(prs,
        "Challenges & Solutions",
        [
            "Challenge: WebM audio unreadable by librosa.",
            "Solution: Client-side decoding and WAV re-encoding via AudioContext.",
            "Challenge: Voice analysis UI hung indefinitely.",
            "Solution: Replaced 'pyin' with 'yin' & added 15-second AbortController.",
            "Challenge: React inputs losing focus.",
            "Solution: Extracted input components outside main render scope.",
            "Challenge: Dragon Breather Mic latency.",
            "Solution: Eagerly initialized mic stream on component mount."
        ],
        notes="Real engineering is solving real bugs. The voiced_flag bug was particularly tricky — it was a silent NameError that caused the backend to crash with no UI feedback, leaving users staring at a spinning loader forever. The solution required both a backend fix and a frontend timeout mechanism.")

    # Slide 15
    add_slide(prs,
        "Future Scope",
        [
            "3–6 Months: Progressive Web App, Real ABDM integration, Trained ML scoring.",
            "6–18 Months: Wearable API integration (SpO2), Telemedicine module, Multi-language.",
            "18+ Months: Hospital backend dashboards, Health Insurance integrations, National ABDM rollout."
        ],
        notes="BreathySync is built for scale. The near-term roadmap focuses on real government integration and replacing simulated ML with trained models. The long-term vision is a platform used by hospitals, insurance companies, and government health programs — making it the backbone of India's digital respiratory health infrastructure.")

    # Slide 16
    add_slide(prs,
        "Conclusion",
        [
            "Non-invasive lung assessment — no hardware, just a phone.",
            "Real-time environmental risk mapping.",
            "Predictive AI that acts 24–48 hours before an attack.",
            "ABHA/UHI government framework integration.",
            "Gamified adherence for sustainable patient engagement.",
            "Thank you!"
        ],
        notes="To summarize: BreathySync transforms a $50,000 clinical problem into a $0 smartphone solution. It's built on sound science, engineered for scale, and designed for India's unique healthcare challenges. Thank you — I'd be happy to take questions or demonstrate the live application.")

    prs.save("BreathySync_Presentation.pptx")
    print("Presentation created successfully as 'BreathySync_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
